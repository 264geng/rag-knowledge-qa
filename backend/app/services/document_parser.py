"""
文档解析服务模块
负责解析 PDF、TXT、Markdown、DOCX、XLSX、PPTX、HTML、CSV 等格式的文档，并进行文本分块
"""

import os
import csv
from typing import List

from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter

from ..config import settings


def parse_pdf(file_path: str) -> str:
    """解析 PDF 文件，提取全部文本内容"""
    reader = PdfReader(file_path)
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)
    return "\n".join(text_parts)


def parse_txt(file_path: str) -> str:
    """解析 TXT 文件"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_markdown(file_path: str) -> str:
    """解析 Markdown 文件（纯文本提取）"""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_docx(file_path: str) -> str:
    """解析 Word 文档，提取所有段落文本"""
    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def parse_xlsx(file_path: str) -> str:
    """解析 Excel 文件，提取所有工作表数据"""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"【工作表: {sheet_name}】")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                text_parts.append(row_text)
    wb.close()
    return "\n".join(text_parts)


def parse_pptx(file_path: str) -> str:
    """解析 PPT 文件，提取所有幻灯片文本"""
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            text_parts.append(f"【幻灯片 {i}】\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts)


def parse_html(file_path: str) -> str:
    """解析 HTML 文件，提取纯文本"""
    from bs4 import BeautifulSoup
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    # 移除 script 和 style 标签
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def parse_csv_file(file_path: str) -> str:
    """解析 CSV 文件"""
    text_parts = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = ", ".join(cell.strip() for cell in row if cell.strip())
            if row_text:
                text_parts.append(row_text)
    return "\n".join(text_parts)


PARSERS = {
    ".pdf": parse_pdf,
    ".txt": parse_txt,
    ".md": parse_markdown,
    ".markdown": parse_markdown,
    ".docx": parse_docx,
    ".xlsx": parse_xlsx,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm": parse_html,
    ".csv": parse_csv_file,
}


def parse_document(file_path: str) -> str:
    """
    根据文件扩展名自动选择解析器
    支持 PDF、TXT、Markdown、DOCX、XLSX、PPTX、HTML、CSV
    """
    ext = os.path.splitext(file_path)[1].lower()
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"不支持的文件格式: {ext}（支持: {', '.join(PARSERS.keys())}）")
    return parser(file_path)


def split_text(text: str) -> List[str]:
    """
    使用 LangChain 的 RecursiveCharacterTextSplitter 对文本进行分块
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.CHUNK_SIZE,
        chunk_overlap=settings.CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""]
    )
    chunks = text_splitter.split_text(text)
    return [chunk for chunk in chunks if chunk.strip()]


def process_document(file_path: str) -> List[str]:
    """
    完整的文档处理流程：
    1. 解析文档提取文本
    2. 将文本分块
    返回分块后的文本列表
    """
    text = parse_document(file_path)
    if not text.strip():
        raise ValueError("文档内容为空")
    chunks = split_text(text)
    return chunks


def get_content_preview(file_path: str, max_length: int = 500) -> str:
    """获取文档内容预览（前 max_length 字符）"""
    text = parse_document(file_path)
    if len(text) > max_length:
        return text[:max_length] + "..."
    return text
